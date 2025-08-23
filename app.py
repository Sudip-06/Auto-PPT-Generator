#!/usr/bin/env python3
"""
Auto PPT Generator — Single-provider only (no env fallbacks), rate-limit aware,
and robust parsing/rendering to keep text inside slides and maximize output.

- Uses ONLY the provider + API key from the form (no env-key fallbacks)
- Gemini 2.5 Pro JSON mode, max_output_tokens=4000; no response.text
- Sectional generation + local MaxiSynth to maximize slide count/details
- Safe margins, autosize, adaptive font, pagination to avoid text overflow
"""

import os
import re
import json
import time
import random
import tempfile
import logging
from io import BytesIO
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple

from flask import Flask, request, jsonify, send_file, render_template
from flask_cors import CORS

# --- LLM SDKs ---
import openai
import anthropic
import google.generativeai as genai

# --- PowerPoint ---
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

# ------------------------------------------------------------------------------
# Logging & App
# ------------------------------------------------------------------------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("app")

app = Flask(__name__)
CORS(app, origins=["*"])
app.config["MAX_CONTENT_LENGTH"] = 16 * 1024 * 1024  # 16MB uploads


# ------------------------------------------------------------------------------
# Simple rate limiter (best-effort per-process)
# ------------------------------------------------------------------------------
class RateLimiter:
    """
    Per-provider best-effort limiter to avoid burst 429s.
    Defaults:
      - google: RPM=2 (Gemini free-tier limit in logs)
      - openai: RPM=30
      - anthropic: RPM=10
    Override via env if needed (e.g., GEMINI_RPM).
    """

    def __init__(self):
        self.last_call_at: Dict[str, float] = {}
        self.rpm = {
            "google": int(os.environ.get("GEMINI_RPM", "2")),
            "openai": int(os.environ.get("OPENAI_RPM", "30")),
            "anthropic": int(os.environ.get("ANTHROPIC_RPM", "10")),
        }

    def before_call(self, provider: str, retry_after: Optional[int] = None):
        now = time.time()
        rpm = self.rpm.get(provider, 10)
        min_gap = 60.0 / max(rpm, 1)
        last = self.last_call_at.get(provider, 0)

        # If server suggested retry_after, prefer it (with jitter)
        if retry_after and retry_after > 0:
            sleep_for = retry_after + random.uniform(0.2, 0.8)
            logger.info(f"[RL] Sleeping {sleep_for:.1f}s due to server retry_after for {provider}")
            time.sleep(sleep_for)
        else:
            gap = (last + min_gap) - now
            if gap > 0:
                sleep_for = gap + random.uniform(0.05, 0.25)
                logger.info(f"[RL] Throttling {provider} for {sleep_for:.2f}s to respect RPM={rpm}")
                time.sleep(sleep_for)
        self.last_call_at[provider] = time.time()


rl = RateLimiter()


# ------------------------------------------------------------------------------
# PPT Generator
# ------------------------------------------------------------------------------
class PPTGenerator:
    def __init__(self):
        self.supported_providers = ["openai", "anthropic", "google"]
        self.max_slides = 12
        self.min_slides = 3
        self._style_ctx: Optional[dict] = None
        self.section_call_cap = 5  # cap sectional calls to be gentle on RPM

    # ========================== PUBLIC ENTRYPOINTS ==========================

    def parse_text_to_slides(
        self, text: str, provider: str, api_key: str, guidance: str = ""
    ) -> List[Dict]:
        """
        Read-first pipeline, using ONLY the chosen provider+key.
        """
        text = (text or "").strip()
        word_count = len(text.split())
        estimated = max(self.min_slides, min(self.max_slides, word_count // 120 + 1))

        # Local deterministic analysis + MaxiSynth plan
        analysis = self._analyze_input_text(text)
        plan = self._maxisynth_plan(analysis, guidance, target_slides=estimated)

        slides: List[Dict] = []
        try:
            logger.info(f"Parsing with {provider}, estimated slides: {estimated}")
            # Primary generation from plan
            response = self._call_llm_with_retry(provider, api_key, self._prompt_from_plan(plan))
            slides = self._robust_json_extraction(response)
            slides = self._validate_and_enhance_slides(slides)

            # If thin, try sectional generation (quota-aware) with same provider
            if len(slides) < max(5, estimated - 1):
                sec = self._sectional_generate(provider, api_key, plan, estimated)
                sec = self._validate_and_enhance_slides(sec)
                if self._score_slides(sec) >= self._score_slides(slides):
                    slides = sec

            # Single refinement pass (same provider)
            try:
                improved = self._refine_with_provider(provider, api_key, json.dumps({"slides": slides}), estimated + 2)
                slides2 = self._robust_json_extraction(improved)
                slides2 = self._validate_and_enhance_slides(slides2)
                if self._score_slides(slides2) >= self._score_slides(slides):
                    slides = slides2
                    logger.info(f"Refinement improved slides to {len(slides)}")
            except Exception:
                logger.warning("Refinement pass failed; keeping initial slides")
        except Exception as e:
            logger.warning(f"Primary generation failed: {e}")

        # If still empty, synthesize locally at max density
        if not slides:
            logger.warning("LLM failed/quota-limited; using MaxiSynth local generation")
            slides = self._maxisynth_local_slides(plan, target=self.max_slides)

        # Ensure conclusion slide
        if not slides or slides[-1].get("slide_type") != "conclusion_slide":
            slides.append(
                {
                    "title": "Conclusion & Next Steps",
                    "content": ["Summary of key points", "Actionable next steps", "Q&A"],
                    "slide_type": "conclusion_slide",
                    "speaker_notes": "Wrap up",
                }
            )

        logger.info(f"Generated {len(slides)} slides successfully")
        return slides[: self.max_slides]

    def create_presentation(self, slides_data: List[Dict], template_file: Optional[str] = None) -> Presentation:
        """Preserve uploaded deck, append generated slides; reuse styles & images."""
        try:
            if template_file:
                prs = Presentation(template_file)
                style = self._extract_style_and_assets(prs)
                logger.info("Using uploaded template/presentation (preserving existing slides, appending new)")
            else:
                prs = Presentation()
                style = {
                    "title_font": None, "body_font": None,
                    "title_color": RGBColor(31, 73, 125), "body_color": RGBColor(64, 64, 64),
                    "images": []
                }
                logger.info("Using default template")
        except Exception as e:
            logger.warning(f"Template error: {e}, using default")
            prs = Presentation()
            style = {
                "title_font": None, "body_font": None,
                "title_color": RGBColor(31, 73, 125), "body_color": RGBColor(64, 64, 64),
                "images": []
            }

        # Pre-render shaping to avoid overflow
        slides_data = self._split_long_bullets(slides_data)
        slides_data = self._paginate_content(slides_data)

        self._style_ctx = style
        for i, s in enumerate(slides_data):
            try:
                self._create_enhanced_slide(prs, s, i)
                if self._style_ctx["images"]:
                    img = self._style_ctx["images"][i % len(self._style_ctx["images"])]
                    try:
                        self._place_logo_safe(prs, prs.slides[-1], img, s.get("slide_type") == "title_slide")
                    except Exception:
                        pass
            except Exception as e:
                logger.error(f"Error creating slide {i}: {e}")
                self._create_fallback_slide(prs, s, i)
        self._style_ctx = None
        return prs

    # ========================== (Helpers) scoring & refine ==========================

    def _score_slides(self, slides: List[Dict]) -> int:
        score = 0
        for s in slides:
            score += min(5, len(s.get("content") or []))
            if s.get("slide_type") == "conclusion_slide":
                score += 1
        return score

    def _refine_with_provider(self, provider: str, api_key: str, slides_json: str, target: int) -> str:
        prompt = f"""
You are refining slides. Improve clarity and actionability. Keep JSON schema identical.
- Total slides ~{target} (±2)
- Each slide: 3–5 concise bullets (<15 words), concrete
- Keep 'slide_type' values; ensure last is 'conclusion_slide'
- Do not add markdown; JSON only
INPUT:
{slides_json}
""".strip()
        return self._call_llm_with_retry(provider, api_key, prompt)

    # ========================== Prompts and local planning ==========================

    def _prompt_from_plan(self, plan: Dict[str, Any]) -> str:
        return f"""
You are a presentation designer.

Use ONLY the PLAN data to build slides. Do not invent facts outside PLAN.
Return ONLY a JSON object (no markdown, no prose).

STRUCTURE:
- 1 title slide for the whole topic.
- {plan.get('target_slides', 8)}±2 total slides.
- Each content slide: 3–5 bullets, each < 15 words, concrete & action-oriented.
- End with a conclusion slide.

PLAN:
{json.dumps(plan, ensure_ascii=False)[:9000]}

JSON FORMAT:
{{
  "presentation_title": "Clear Title",
  "slides": [
    {{"title":"Welcome","content":["What this covers","Key value"],"slide_type":"title_slide","speaker_notes":"short"}},
    {{"title":"Main Point","content":["Specific fact","Example","Action"],"slide_type":"content_slide","speaker_notes":"short"}},
    {{"title":"Conclusion & Next Steps","content":["Summary","Next steps","Q&A"],"slide_type":"conclusion_slide","speaker_notes":"short"}}
  ]
}}
""".strip()

    def _analyze_input_text(self, text: str) -> Dict[str, Any]:
        t = text.replace("_x000D_", " ").replace("\r\n", "\n").replace("\r", "\n")

        title_match = re.search(r"^\s*#\s+(.+)$", t, re.MULTILINE)
        if title_match:
            title = title_match.group(1).strip()
        else:
            first_line = next((ln.strip() for ln in t.split("\n") if ln.strip()), "")
            title = " ".join(first_line.split()[:12]) or "Presentation"

        section_headers = [(m.start(), m.group(0), m.group(1).strip())
                           for m in re.finditer(r"^\s*(##+)\s+(.+)$", t, re.MULTILINE)]
        sections: List[Dict[str, Any]] = []
        if section_headers:
            for idx, (pos, hashes, heading) in enumerate(section_headers):
                start = pos + len(hashes) + 1 + len(heading)
                end = section_headers[idx + 1][0] if idx + 1 < len(section_headers) else len(t)
                body = t[start:end].strip()
                sections.append({"heading": heading[:80], "text": body})
        else:
            paras = [p.strip() for p in re.split(r"\n{2,}", t) if p.strip()]
            for i, chunk in enumerate(paras[:12]):
                sections.append({"heading": f"Section {i+1}", "text": chunk})

        facts, nums, dates = self._harvest_facts(t)
        return {
            "title": title,
            "sections": sections[:10],
            "facts": facts[:30],
            "numbers": nums[:30],
            "dates": dates[:20],
        }

    def _harvest_facts(self, text: str) -> Tuple[List[str], List[str], List[str]]:
        nums = re.findall(r"\b(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?%?\b", text)
        date_patterns = [
            r"\b(?:Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|"
            r"Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)[ ,.-]?\s?\d{1,2},?\s?\d{2,4}\b",
            r"\b\d{4}-\d{2}-\d{2}\b",
            r"\b(?:Q[1-4]\s*\d{4})\b",
            r"\b\d{2}/\d{2}/\d{4}\b",
        ]
        dates = []
        for dp in date_patterns:
            dates += re.findall(dp, text, flags=re.IGNORECASE)
        dates = list(dict.fromkeys(dates))

        sents = re.split(r"(?<=[.!?])\s+", text)
        facts = []
        for s in sents:
            if len(s) < 30:
                continue
            if re.search(r"\d", s) or re.search(r"\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+){0,3}\b", s):
                facts.append(s.strip())

        return facts, nums, dates

    def _maxisynth_plan(self, analysis: Dict[str, Any], guidance: str, target_slides: int) -> Dict[str, Any]:
        title = analysis.get("title") or "Presentation"
        secs = analysis.get("sections", [])
        facts = analysis.get("facts", [])
        nums = analysis.get("numbers", [])
        dates = analysis.get("dates", [])
        style = guidance or "professional, clear, and engaging presentation"

        sec_payload = []
        for s in secs:
            lines = [ln.strip() for ln in s.get("text","").split("\n") if ln.strip()]
            bullets = [re.sub(r"^(\*|-|•|\d+\.)\s*", "", ln).strip() for ln in lines if re.match(r"^(\*|-|•|\d+\.)\s+", ln)]
            if not bullets:
                sentences = re.split(r"(?<=[.!?])\s+", s.get("text",""))
                bullets = [x.strip() for x in sentences if len(x.strip()) > 30][:5]
            sec_payload.append({"title": s["heading"], "bullets": bullets[:6]})

        key_metrics = list(dict.fromkeys(nums))[:12]
        key_dates = dates[:10]
        key_facts = facts[:12]

        return {
            "topic": title,
            "style": style,
            "sections": sec_payload,
            "key_metrics": key_metrics,
            "key_dates": key_dates,
            "key_facts": key_facts,
            "target_slides": min(self.max_slides, max(target_slides, 8)),
        }

    def _maxisynth_local_slides(self, plan: Dict[str, Any], target: int) -> List[Dict]:
        slides: List[Dict] = []
        topic = plan.get("topic", "Presentation")

        slides.append({
            "title": topic,
            "content": ["Overview", "Scope & objectives", "What you’ll learn"],
            "slide_type": "title_slide",
            "speaker_notes": f"Introduce {topic}"
        })

        for sec in plan.get("sections", [])[:6]:
            bullets = sec.get("bullets") or []
            bullets = [b[:150] for b in bullets if b.strip()]
            if not bullets:
                continue
            slides.append({
                "title": sec.get("title", "Section"),
                "content": bullets[:5],
                "slide_type": "content_slide",
                "speaker_notes": f"Discuss {sec.get('title','section')}"
            })
            if len(slides) >= target - 4:
                break

        facts = plan.get("key_facts", [])
        if facts:
            slides.append({"title": "Key Facts", "content": [f[:150] for f in facts[:5]],
                           "slide_type": "content_slide", "speaker_notes": "Highlight key facts"})

        metrics = plan.get("key_metrics", [])
        if metrics:
            slides.append({"title": "Metrics & Figures", "content": [m[:60] for m in metrics[:5]],
                           "slide_type": "content_slide", "speaker_notes": "Explain the numbers"})

        if plan.get("key_dates"):
            slides.append({"title": "Timeline", "content": plan["key_dates"][:5],
                           "slide_type": "content_slide", "speaker_notes": "Walk through dates"})

        faqs = self._derive_faqs(plan)
        if faqs:
            slides.append({"title": "FAQs", "content": faqs[:5], "slide_type": "content_slide",
                           "speaker_notes": "Answer common questions"})

        risks = self._derive_risks(plan)
        if risks:
            slides.append({"title": "Risks & Mitigations", "content": risks[:5], "slide_type": "content_slide",
                           "speaker_notes": "Call out risks"})

        slides.append({"title": "Next Steps",
                       "content": ["Assign owners", "Set milestones", "Share deck", "Kickoff meeting"],
                       "slide_type": "content_slide", "speaker_notes": "Action plan"})
        return slides

    def _derive_faqs(self, plan: Dict[str, Any]) -> List[str]:
        faqs = []
        for f in plan.get("key_facts", []):
            if "?" in f and len(faqs) < 5:
                faqs.append(f)
        if len(faqs) < 3:
            faqs.extend(["Who is the target audience?",
                         "What is the expected outcome?",
                         "What are the success metrics?"])
        return faqs

    def _derive_risks(self, plan: Dict[str, Any]) -> List[str]:
        risks = []
        text_blobs = " ".join([" ".join(sec.get("bullets", [])) for sec in plan.get("sections", [])])
        if re.search(r"\brisks?\b|\bchallenge(s)?\b|\bissue(s)?\b", text_blobs, re.IGNORECASE):
            risks.append("Scope creep without clear milestones")
            risks.append("Missing data or unclear ownership")
        risks.append("Timeline slippage due to dependencies")
        risks.append("Stakeholder misalignment")
        risks.append("Insufficient testing before launch")
        return risks

    # ========================== LLM CALLS (with rate limiting) ==========================

    def _call_llm_with_retry(self, provider: str, api_key: str, prompt: str, max_retries: int = 2) -> str:
        last_err = None
        for attempt in range(max_retries + 1):
            try:
                retry_after = None
                if provider == "google":
                    rl.before_call("google")
                elif provider == "anthropic":
                    rl.before_call("anthropic")
                elif provider == "openai":
                    rl.before_call("openai")

                if provider == "openai":
                    return self._call_openai(prompt, api_key)
                if provider == "anthropic":
                    return self._call_anthropic(prompt, api_key)
                if provider == "google":
                    return self._call_gemini(prompt, api_key)

                raise ValueError(f"Unsupported provider: {provider}")
            except Exception as e:
                last_err = e
                msg = str(e)
                # server-provided retry hint like "retry_delay { seconds: 58 }"
                m = re.search(r"retry[_ ]delay[^0-9]*([0-9]+)", msg, flags=re.IGNORECASE)
                retry_after = int(m.group(1)) if m else None

                if attempt < max_retries:
                    logger.warning(f"Attempt {attempt + 1} failed: {e}, retrying...")
                    if provider in ("google", "openai", "anthropic"):
                        rl.before_call(provider, retry_after=retry_after)
                else:
                    break
        raise last_err or RuntimeError("LLM call failed")

    def _call_openai(self, prompt: str, api_key: str) -> str:
        openai.api_key = api_key
        try:
            resp = openai.ChatCompletion.create(
                model="gpt-5",
                messages=[
                    {"role": "system", "content": "Output valid JSON only. No markdown."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.15,
                max_tokens=1800,
            )
            content = (resp.choices[0].message.content or "").strip()
            if not content:
                raise ValueError("Empty response from OpenAI")
            return content
        except Exception as e:
            msg = str(e).lower()
            if "api key" in msg or "auth" in msg:
                raise ValueError("Invalid OpenAI API key")
            if "rate" in msg or "quota" in msg or "limit" in msg:
                raise ValueError("OpenAI rate limit exceeded")
            raise

    def _call_anthropic(self, prompt: str, api_key: str) -> str:
        try:
            client = anthropic.Anthropic(api_key=api_key)
            resp = client.messages.create(
                model="claude-opus-4-1",
                max_tokens=1800,
                temperature=0.15,
                system="Return ONLY valid JSON (no markdown).",
                messages=[{"role": "user", "content": prompt}],
            )
            parts: List[str] = []
            for block in (resp.content or []):
                if getattr(block, "type", "") == "text" and getattr(block, "text", ""):
                    parts.append(block.text)
            text = ("\n".join(parts)).strip()
            if not text:
                raise ValueError("Empty response from Claude")
            return text
        except anthropic.AuthenticationError:
            raise ValueError("Invalid Anthropic API key")
        except anthropic.RateLimitError:
            raise ValueError("Anthropic rate limit exceeded")
        except Exception:
            raise

    def _call_gemini(self, prompt: str, api_key: str) -> str:
        """
        Gemini 2.5 Pro, JSON mime mode; avoid response.text; collect parts.
        """
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel("gemini-2.5-pro")
            gen_cfg = {
                "temperature": 0.1,
                "max_output_tokens": 4000,   # per your request
                "candidate_count": 1,
                "response_mime_type": "application/json",
                "top_p": 0.8,
                "top_k": 40,
            }
            resp = model.generate_content(prompt, generation_config=gen_cfg)
            texts = []
            for cand in getattr(resp, "candidates", []) or []:
                content = getattr(cand, "content", None)
                for part in (getattr(content, "parts", None) or []):
                    t = getattr(part, "text", None)
                    if t:
                        texts.append(t)
            out = ("\n".join(texts)).strip()
            if not out:
                raise ValueError("Empty response from Google")
            return out
        except Exception as e:
            msg = str(e)
            if "api key" in msg.lower():
                raise ValueError("Invalid Google API key")
            if "429" in msg or "rate limit" in msg.lower() or "quota" in msg.lower():
                raise ValueError("Google rate limit exceeded")
            if "500" in msg:
                raise ValueError("Google error: transient server issue (500). Please retry.")
            raise

    # ======= Sectional generation (same provider only) =======

    def _sectional_generate(self, provider: str, api_key: str, plan: Dict[str, Any], target_slides: int) -> List[Dict]:
        sections = plan.get("sections", [])[: self.section_call_cap]
        slides: List[Dict] = []
        topic = plan.get("topic") or "Presentation"

        slides.append({
            "title": topic,
            "content": [f"Overview of {topic}", "Key takeaways", "Agenda"],
            "slide_type": "title_slide",
            "speaker_notes": f"Introduce {topic}"
        })

        per_section = max(1, min(2, target_slides // max(1, len(sections))))

        for sec in sections:
            sec_prompt = f"""
Build JSON slides ONLY for the SECTION, grounded strictly in its bullets.
No markdown, no prose. JSON only.

OUTPUT: {per_section} content slide(s), each with 3–5 bullets (<15 words).
SECTION:
{json.dumps(sec, ensure_ascii=False)[:2500]}

JSON FORMAT:
{{"slides":[{{"title":"...","content":["..."],"slide_type":"content_slide","speaker_notes":"short"}}]}}
""".strip()
            try:
                if provider == "google":
                    rl.before_call("google")
                elif provider == "anthropic":
                    rl.before_call("anthropic")
                elif provider == "openai":
                    rl.before_call("openai")

                if provider == "openai":
                    out = self._call_openai(sec_prompt, api_key)
                elif provider == "anthropic":
                    out = self._call_anthropic(sec_prompt, api_key)
                else:
                    out = self._call_gemini(sec_prompt, api_key)

                sec_slides = self._robust_json_extraction(out)
                for ss in sec_slides:
                    if ss.get("slide_type") == "title_slide":
                        ss["slide_type"] = "content_slide"
                slides.extend(sec_slides)
                if len(slides) >= self.max_slides - 1:
                    break
            except Exception as e:
                logger.warning(f"Sectional gen error: {e}; continuing.")
                if "rate limit" in str(e).lower() or "quota" in str(e).lower():
                    break  # stop on quota to avoid thrash

        slides.append({
            "title": "Conclusion & Next Steps",
            "content": ["Summary of key points", "Actionable next steps", "Q&A"],
            "slide_type": "conclusion_slide",
            "speaker_notes": "Wrap up"
        })
        return slides

    # ========================== JSON EXTRACTION & VALIDATION ==========================

    def _robust_json_extraction(self, response: str) -> List[Dict]:
        cleaned = (response or "").strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.strip("`").strip()
            if cleaned.lower().startswith("json"):
                cleaned = cleaned[4:].strip()

        match = re.search(r"\{.*\"slides\".*\}[\s\S]*$", cleaned)
        candidate = match.group(0) if match else cleaned

        try:
            data = json.loads(candidate)
        except Exception:
            arr = re.search(r"\"slides\"\s*:\s*(\[[\s\S]*?\])", cleaned)
            if arr:
                try:
                    slides = json.loads(arr.group(1))
                    return self._validate_and_enhance_slides(slides)
                except Exception:
                    pass
            return self._md_to_slides(cleaned)

        if not (isinstance(data, dict) and isinstance(data.get("slides"), list) and data["slides"]):
            return self._md_to_slides(cleaned)

        return self._validate_and_enhance_slides(data)

    def _validate_and_enhance_slides(self, data_or_slides) -> List[Dict]:
        slides = data_or_slides if isinstance(data_or_slides, list) else data_or_slides.get("slides", [])
        out: List[Dict] = []

        for i, s in enumerate(slides):
            title = self._clean_text(str(s.get("title", f"Slide {i+1}"))) or f"Slide {i+1}"
            notes = self._clean_text(str(s.get("speaker_notes", ""))) or f"Notes for {title}"
            stype = s.get("slide_type", "content_slide")

            content = s.get("content", [])
            if isinstance(content, str):
                content = [content]
            content = [self._clean_text(str(x))[:150] for x in content if str(x).strip()]

            if len(content) < 3 and i > 0:
                content += [f"Key insight about {title.lower()}"]
            content = [c for c in content if c][:5]
            if not content and i > 0:
                content = [f"Overview of {title}", "Main takeaway", "Next step"]

            if stype not in {"title_slide", "content_slide", "conclusion_slide", "section_header"}:
                stype = "content_slide"
            if i == 0:
                stype = "title_slide"

            out.append({"title": title, "content": content, "slide_type": stype, "speaker_notes": notes})

        return out[: self.max_slides]

    def _md_to_slides(self, text: str) -> List[Dict]:
        lines = (text or "").splitlines()
        slides: List[Dict] = []
        cur: Optional[Dict] = None

        def push():
            nonlocal cur
            if not cur:
                return
            bullets = [re.sub(r"\s+", " ", b).strip() for b in cur.get("content", []) if b.strip()]
            bullets = [b[:150] for b in bullets] or [f"Overview of {cur['title']}"]
            if len(bullets) < 3:
                bullets += [" "] * (3 - len(bullets))
            cur["content"] = bullets[:5]
            slides.append(cur)
            cur = None

        for raw in lines:
            ln = raw.strip()
            if not ln:
                continue
            if ln.startswith("# "):
                if cur:
                    push()
                t = ln[2:].strip() or "Presentation"
                cur = {"title": t, "content": [], "slide_type": "title_slide", "speaker_notes": f"Intro: {t}"}
            elif ln.startswith("## ") or ln.startswith("### "):
                if cur:
                    push()
                t = ln.split(" ", 1)[1].strip() or "Section"
                cur = {"title": t, "content": [], "slide_type": "content_slide", "speaker_notes": f"Discuss: {t}"}
            elif re.match(r"^(\*|-|•|\d+\.)\s+", ln):
                if not cur:
                    cur = {"title": "Key Points", "content": [], "slide_type": "content_slide", "speaker_notes": "Key points"}
                b = re.sub(r"^(\*|-|•|\d+\.)\s+", "", ln).strip()
                if b:
                    cur["content"].append(b)
            else:
                if not cur:
                    cur = {"title": "Overview", "content": [], "slide_type": "content_slide", "speaker_notes": "Overview"}
                if len(ln) > 40:
                    cur["content"].append(ln)

        if cur:
            push()

        if not slides:
            return self._create_default_slides(text)
        if not any(s.get("slide_type") == "conclusion_slide" for s in slides[-2:]):
            slides.append({
                "title": "Conclusion & Next Steps",
                "content": ["Summary of key points", "Actionable next steps", "Q&A"],
                "slide_type": "conclusion_slide",
                "speaker_notes": "Wrap up",
            })
        return slides

    # ========================== CLEANING & DEFAULTS ==========================

    def _clean_text(self, text: str) -> str:
        if not text:
            return ""
        text = text.replace("_x000D_", " ").replace("\r\n", " ").replace("\r", " ")
        cleaned = re.sub(r"\s+", " ", text.strip())
        cleaned = re.sub(r"[^\w\s\-.,!?()&%$#@:/]", "", cleaned)
        return cleaned

    def _create_default_slides(self, original_text: str) -> List[Dict]:
        words = (original_text or "").split()
        chunk = " ".join(words[:250])
        slides = [
            {"title": "Presentation Overview",
             "content": ["Key insights from your content", "Structured information", "Professional presentation"],
             "slide_type": "title_slide", "speaker_notes": "Intro"}
        ]
        sentences = re.split(r"[.!?]+", chunk)
        bullets = [s.strip() for s in sentences if len(s.strip()) > 20][:4]
        if bullets:
            slides.append({"title": "Key Points", "content": bullets, "slide_type": "content_slide", "speaker_notes": "Explain"})
        slides.append({"title": "Conclusion & Next Steps", "content": ["Summary", "Next steps", "Q&A"], "slide_type": "conclusion_slide", "speaker_notes": "Wrap up"})
        return slides

    # ========================== LAYOUT / RENDERING ==========================

    def _split_long_bullets(self, slides: List[Dict], max_len: int = 140) -> List[Dict]:
        out = []
        for s in slides:
            bullets = []
            for b in (s.get("content") or []):
                t = (b or "").strip()
                if len(t) <= max_len:
                    bullets.append(t)
                else:
                    split_pt = max(
                        t.rfind(". ", 0, max_len),
                        t.rfind("; ", 0, max_len),
                        t.rfind(", ", 0, max_len),
                        t.rfind(" - ", 0, max_len),
                    )
                    if split_pt < 60:
                        split_pt = max_len
                    bullets.append(t[:split_pt].strip())
                    rest = t[split_pt:].lstrip(".;,- ").strip()
                    if rest:
                        bullets.append(rest[:max_len].strip())
            s2 = dict(s)
            s2["content"] = bullets[:5] if bullets else s.get("content", [])
            out.append(s2)
        return out

    def _paginate_content(self, slides: List[Dict]) -> List[Dict]:
        paginated: List[Dict] = []
        limit = 600  # rough capacity for 18–20pt with bullets
        for s in slides:
            bullets = s.get("content") or []
            total = sum(len(b or "") for b in bullets)
            if total <= limit or s.get("slide_type") == "title_slide":
                paginated.append(s)
                continue
            chunk, clen, idx = [], 0, 1
            for b in bullets:
                bl = len(b or "")
                if chunk and (clen + bl) > limit:
                    paginated.append({"title": f"{s.get('title','Slide')} (cont. {idx})",
                                      "content": chunk[:5], "slide_type": "content_slide",
                                      "speaker_notes": s.get("speaker_notes","")})
                    chunk, clen, idx = [], 0, idx + 1
                chunk.append(b); clen += bl
            if chunk:
                suf = "" if idx == 1 else f" (cont. {idx})"
                paginated.append({"title": f"{s.get('title','Slide')}{suf}",
                                  "content": chunk[:5], "slide_type": "content_slide",
                                  "speaker_notes": s.get("speaker_notes","")})
        return paginated

    def _safe_rect(self, prs: Presentation) -> Tuple[Emu, Emu, Emu, Emu]:
        sw, sh = prs.slide_width, prs.slide_height
        margin_x = Emu(Inches(0.8))
        top = Emu(Inches(1.6))
        bottom = Emu(Inches(0.8))
        left = margin_x
        width = sw - Emu(Inches(0.8)) - margin_x
        height = sh - top - bottom
        return left, top, width, height

    def _apply_textframe_defaults(self, tf, base_pt=20):
        try: tf.word_wrap = True
        except: pass
        try: tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except: pass
        try:
            tf.margin_left = Emu(Inches(0.1))
            tf.margin_right = Emu(Inches(0.1))
            tf.margin_top = Emu(Inches(0.05))
            tf.margin_bottom = Emu(Inches(0.05))
        except: pass
        size = base_pt
        joined = " ".join(p.text for p in tf.paragraphs if getattr(p, "text", ""))
        L = len(joined)
        size = 16 if L > 500 else (18 if L > 350 else base_pt)
        try:
            p0 = tf.paragraphs[0]
            p0.font.size = Pt(size)
            ctx = getattr(self, "_style_ctx", None) or {}
            if ctx.get("body_font"):
                p0.font.name = ctx["body_font"]
            p0.font.color.rgb = (ctx.get("body_color", RGBColor(64, 64, 64)))
        except: pass
        return size

    def _ideal_bullet_font_size(self, bullets: List[str]) -> int:
        n = len(bullets or [])
        avg = sum(len(b or "") for b in bullets) / max(n, 1)
        size = 20
        if n >= 5 or avg > 90: size = 18
        if n >= 6 or avg > 120: size = 16
        return size

    def _place_logo_safe(self, prs: Presentation, slide, img_path: str, is_title: bool):
        sw, sh = prs.slide_width, prs.slide_height
        h_in = 0.9 if is_title else 1.0
        pic = slide.shapes.add_picture(img_path, Emu(0), Emu(0), height=Emu(Inches(h_in)))
        margin = Emu(Inches(0.3 if is_title else 0.4))
        pic.left = sw - margin - pic.width
        pic.top = (Emu(Inches(0.3)) if is_title else (sh - margin - pic.height))

    def _extract_style_and_assets(self, prs: Presentation) -> dict:
        style = {
            "title_font": None, "body_font": None,
            "title_color": RGBColor(31, 73, 125), "body_color": RGBColor(64, 64, 64),
            "images": [],
        }
        try:
            if len(prs.slides):
                s0 = prs.slides[0]
                if getattr(s0.shapes, "title", None) and s0.shapes.title.has_text_frame:
                    p = s0.shapes.title.text_frame.paragraphs[0].font
                    if getattr(p, "name", None): style["title_font"] = p.name
                    if getattr(p.color, "rgb", None): style["title_color"] = p.color.rgb
                for ph in s0.placeholders:
                    if hasattr(ph, "text_frame") and ph is not getattr(s0.shapes, "title", None):
                        pf = ph.text_frame.paragraphs[0].font
                        if getattr(pf, "name", None): style["body_font"] = style["body_font"] or pf.name
                        if getattr(pf.color, "rgb", None): style["body_color"] = pf.color.rgb
                        break
        except: pass
        try:
            for s in prs.slides:
                for shp in s.shapes:
                    if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            blob = shp.image.blob
                            ext = shp.image.ext or "png"
                            tf = tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}")
                            tf.write(blob); tf.flush(); tf.close()
                            style["images"].append(tf.name)
                        except: continue
        except: pass
        return style

    def _create_enhanced_slide(self, prs: Presentation, s: Dict, idx: int):
        if s.get("slide_type") == "title_slide" and idx == 0:
            layout_idx = 0 if len(prs.slide_layouts) > 0 else 1
        elif s.get("content"):
            layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        else:
            layout_idx = 5 if len(prs.slide_layouts) > 5 else 1
        try: lay = prs.slide_layouts[layout_idx]
        except: lay = prs.slide_layouts[0]
        slide = prs.slides.add_slide(lay)

        # Title
        try:
            title_text = s["title"]
            if getattr(slide.shapes, "title", None):
                slide.shapes.title.text = title_text
                self._style_title(slide.shapes.title, idx == 0)
            else:
                sw = prs.slide_width
                left = Emu(Inches(0.8)); top = Emu(Inches(0.5))
                width = sw - left - Emu(Inches(0.8)); height = Emu(Inches(1.2))
                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame; tf.clear()
                p = tf.paragraphs[0]; p.text = title_text
                p.font.bold = True; p.font.size = Pt(36 if idx == 0 else 30)
                ctx = getattr(self, "_style_ctx", None) or {}
                if ctx.get("title_font"): p.font.name = ctx["title_font"]
                try: p.font.color.rgb = ctx.get("title_color", RGBColor(31, 73, 125))
                except: pass
        except Exception as e:
            logger.debug(f"Title add failed: {e}")

        # Content
        try:
            bullets = s.get("content") or []
            if bullets:
                ph = None
                for ph_i in getattr(slide, "placeholders", []):
                    if hasattr(ph_i, "text_frame") and ph_i != getattr(slide.shapes, "title", None):
                        ph = ph_i; break
                use_ph = False
                if ph and hasattr(ph, "height") and hasattr(ph, "width"):
                    try:
                        use_ph = (ph.height >= Emu(Inches(2.0)) and ph.width >= Emu(Inches(5.0)))
                    except Exception:
                        use_ph = True
                if use_ph:
                    tf = ph.text_frame; tf.clear()
                    base = self._ideal_bullet_font_size(bullets)
                    self._apply_textframe_defaults(tf, base_pt=base)
                    for i, t in enumerate(bullets):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = t; p.level = 0
                        try:
                            p.font.size = Pt(base); p.space_after = Pt(6)
                            ctx = getattr(self, "_style_ctx", None) or {}
                            if ctx.get("body_font"): p.font.name = ctx["body_font"]
                            p.font.color.rgb = ctx.get("body_color", RGBColor(64, 64, 64))
                        except: pass
                else:
                    left, top, width, height = self._safe_rect(prs)
                    tb = slide.shapes.add_textbox(left, top, width, height)
                    tf = tb.text_frame; tf.clear()
                    base = self._ideal_bullet_font_size(bullets)
                    self._apply_textframe_defaults(tf, base_pt=base)
                    for i, t in enumerate(bullets):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = t; p.level = 0
                        try:
                            p.font.size = Pt(base); p.space_after = Pt(6)
                            ctx = getattr(self, "_style_ctx", None) or {}
                            if ctx.get("body_font"): p.font.name = ctx["body_font"]
                            p.font.color.rgb = ctx.get("body_color", RGBColor(64, 64, 64))
                        except: pass
        except Exception as e:
            logger.warning(f"Could not add content: {e}")
            self._add_basic_content(prs, slide, s.get("content", []))

        # Speaker notes
        self._add_speaker_notes(slide, s.get("speaker_notes", ""))

    def _style_title(self, title_shape, is_main: bool):
        try:
            if title_shape.text_frame:
                p = title_shape.text_frame.paragraphs[0]; f = p.font
                f.size = Pt(36 if is_main else 30); f.bold = True
                ctx = getattr(self, "_style_ctx", None) or {}
                if ctx.get("title_font"): f.name = ctx["title_font"]
                try: f.color.rgb = ctx.get("title_color", RGBColor(31, 73, 125))
                except: pass
        except: pass

    def _add_basic_content(self, prs: Presentation, slide, content_list: List[str]):
        try:
            left, top, width, height = self._safe_rect(prs)
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame; tf.clear()
            base = self._ideal_bullet_font_size(content_list)
            self._apply_textframe_defaults(tf, base_pt=base)
            for i, t in enumerate(content_list):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = t
                try:
                    p.font.size = Pt(base)
                    ctx = getattr(self, "_style_ctx", None) or {}
                    if ctx.get("body_font"): p.font.name = ctx["body_font"]
                    p.font.color.rgb = ctx.get("body_color", RGBColor(64, 64, 64))
                except: pass
        except Exception as e:
            logger.warning(f"Basic content addition failed: {e}")

    def _add_speaker_notes(self, slide, notes: str):
        if not notes: return
        try:
            ns = slide.notes_slide
            ns.notes_text_frame.text = notes
        except: pass

    def _create_fallback_slide(self, prs: Presentation, s: Dict, idx: int):
        try:
            lay = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(lay)
            sw = prs.slide_width
            left = Emu(Inches(0.8)); top = Emu(Inches(0.5))
            width = sw - left - Emu(Inches(0.8))
            tb = slide.shapes.add_textbox(left, top, width, Emu(Inches(1.2)))
            tf = tb.text_frame; tf.text = s.get("title", f"Slide {idx+1}")
            tf.paragraphs[0].font.size = Pt(32); tf.paragraphs[0].font.bold = True
            self._add_basic_content(prs, slide, s.get("content") or [])
        except Exception as e:
            logger.error(f"Even fallback slide creation failed: {e}")


# ------------------------------------------------------------------------------
# Initialize generator
# ------------------------------------------------------------------------------
ppt_generator = PPTGenerator()


# ------------------------------------------------------------------------------
# Routes
# ------------------------------------------------------------------------------
@app.route("/")
def index():
    try:
        return render_template("index.html")  # unchanged
    except Exception:
        return "<h1>Auto PPT Generator</h1><p>Server running but template missing.</p>", 200


@app.route("/generate", methods=["POST"])
def generate_presentation():
    try:
        if "input_text" not in request.form:
            return jsonify({"error": "No input text provided"}), 400

        input_text = (request.form["input_text"] or "").strip()
        guidance = (request.form.get("guidance", "") or "").strip()
        provider = (request.form.get("provider", "openai") or "").strip().lower()
        api_key = (request.form.get("api_key", "") or "").strip()

        if len(input_text) < 20:
            return jsonify({"error": "Please provide at least 20 characters of content"}), 400
        if not api_key:
            return jsonify({"error": "API key is required"}), 400
        if provider not in ppt_generator.supported_providers:
            return jsonify({"error": f"Unsupported provider: {provider}"}), 400

        template_path = None
        if "template_file" in request.files:
            tf = request.files["template_file"]
            if tf and tf.filename and tf.filename.lower().endswith((".pptx", ".potx")):
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    tf.save(tmp.name)
                    template_path = tmp.name

        try:
            slides_data = ppt_generator.parse_text_to_slides(
                text=input_text, provider=provider, api_key=api_key, guidance=guidance
            )
            if not slides_data:
                return jsonify({"error": "Could not generate slides from content"}), 400

            prs = ppt_generator.create_presentation(slides_data, template_path)

            bio = BytesIO()
            prs.save(bio); bio.seek(0)
            return send_file(
                bio,
                as_attachment=True,
                download_name=f'presentation_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx',
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except ValueError as e:
            return jsonify({"error": str(e)}), 400
        except Exception as e:
            logger.error(f"Generation error: {e}")
            return jsonify({"error": "Presentation generation failed. Please try again."}), 500
        finally:
            if template_path and os.path.exists(template_path):
                try: os.unlink(template_path)
                except: pass
    except Exception as e:
        logger.error(f"Request error: {e}")
        return jsonify({"error": "Server error"}), 500


@app.route("/health")
def health():
    return jsonify({"status": "healthy", "timestamp": datetime.utcnow().isoformat()})


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)
